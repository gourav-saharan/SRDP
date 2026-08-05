import customtkinter as ctk
import tkinter as tk
from tkinter import filedialog, messagebox
import os
import sys
import tempfile
from datetime import datetime
import matplotlib

matplotlib.use("TkAgg", force=True)
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg, NavigationToolbar2Tk
import pandas as pd
import numpy as np
from PIL import Image

try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.util import Cm, Emu
    has_pptx = True
except ImportError:
    has_pptx = False

try:
    import scipy.signal as signal
    import scipy.ndimage as ndimage
    has_scipy = True
except ImportError:
    has_scipy = False

def apply_lpf(data, cutoff=5, fs=100, order=4):
    if not has_scipy: return data
    try:
        nyq = 0.5 * fs
        normal_cutoff = cutoff / nyq
        b, a = signal.butter(order, normal_cutoff, btype='low', analog=False)
        return signal.filtfilt(b, a, data)
    except: return data

def apply_elliptic(data, cutoff=5, fs=100, order=5, rp=1.0, rs=40.0):
    if not has_scipy: return data
    try:
        nyq = 0.5 * fs
        normal_cutoff = cutoff / nyq
        if normal_cutoff <= 0 or normal_cutoff >= 1: return data
        b, a = signal.ellip(order, rp, rs, normal_cutoff, btype='low', analog=False)
        return signal.filtfilt(b, a, data)
    except: return data

def apply_savgol(data, window=51, polyorder=3):
    if not has_scipy: return data
    try:
        w = int(window)
        if w % 2 == 0: w += 1
        if w < polyorder + 2: w = polyorder + 2
        if len(data) < w: return data
        return signal.savgol_filter(data, w, polyorder)
    except: return data

def apply_median(data, window=5):
    if not has_scipy: return data
    try:
        w = int(window)
        if w % 2 == 0: w += 1
        return signal.medfilt(data, w)
    except: return data

def apply_kalman(data, process_var=1e-5, measurement_var=0.1):
    try:
        n = len(data)
        posteri_estimate = np.zeros(n)
        posteri_error_estimate = np.zeros(n)
        
        posteri_estimate[0] = data[0] if not np.isnan(data[0]) else 0.0
        posteri_error_estimate[0] = 1.0
        
        for i in range(1, n):
            priori_estimate = posteri_estimate[i-1]
            priori_error_estimate = posteri_error_estimate[i-1] + process_var
            
            blending_factor = priori_error_estimate / (priori_error_estimate + measurement_var)
            val = data[i] if not np.isnan(data[i]) else priori_estimate
            posteri_estimate[i] = priori_estimate + blending_factor * (val - priori_estimate)
            posteri_error_estimate[i] = (1 - blending_factor) * priori_error_estimate
            
        return posteri_estimate
    except: return data

def apply_notch(data, freq=50, fs=100, Q=30):
    if not has_scipy: return data
    try:
        nyq = 0.5 * fs
        w0 = freq / nyq
        if w0 <= 0 or w0 >= 1: return data
        b, a = signal.iirnotch(w0, Q)
        return signal.filtfilt(b, a, data)
    except: return data

def apply_hampel(data, window=5, n_sigmas=3):
    try:
        w = int(window)
        if w % 2 == 0: w += 1
        s = pd.Series(data)
        rolling_median = s.rolling(window=w, center=True).median()
        MAD = s.rolling(window=w, center=True).apply(lambda x: np.median(np.abs(x - np.median(x))), raw=True)
        threshold = n_sigmas * 1.4826 * MAD
        
        outlier_idx = np.abs(s - rolling_median) > threshold
        new_data = data.copy()
        new_data[outlier_idx] = rolling_median[outlier_idx]
        
        if np.isnan(new_data).any():
            new_data = pd.Series(new_data).interpolate(limit_direction='both').values
        return new_data
    except: return data

def apply_gaussian(data, sigma=2):
    if not has_scipy: return data
    try:
        return ndimage.gaussian_filter1d(data, sigma=max(1.0, float(sigma)))
    except: return data

def apply_wiener(data, window=5):
    if not has_scipy: return data
    try:
        w = int(window)
        return signal.wiener(data, mysize=w)
    except: return data

def apply_wavelet(data, level=1):
    try:
        import pywt
        wavelet = 'db4'
        coeff = pywt.wavedec(data, wavelet, mode="symmetric")
        sigma = (1/0.6745) * np.median(np.abs(coeff[-1] - np.median(coeff[-1])))
        uthresh = sigma * np.sqrt(2*np.log(len(data)))
        
        coeff[1:] = [pywt.threshold(i, value=uthresh, mode='soft') for i in coeff[1:]]
        rec = pywt.waverec(coeff, wavelet, mode="symmetric")
        return rec[:len(data)]
    except: return data

class GraphFrame(ctk.CTkFrame):
    def __init__(self, parent, controller):
        super().__init__(parent)
        self.controller = controller
        
        self.grid_rowconfigure(4, weight=1)
        self.grid_columnconfigure(0, weight=1)
        self.y_axis_limit_values = {}
        self.y_axis_limit_entries = {}
        self.y_axes = []
        self.extra_axes = []
        self.ppt_slide_entries = []
        self.ppt_preview_images = []
        self.ppt_output_dir = os.path.join(tempfile.gettempdir(), "SRDP_PPT_Graphs")
        os.makedirs(self.ppt_output_dir, exist_ok=True)
        self.ppt_template_path = self.resource_path("Temp.pptx")
        self.ppt_graph_box = None
        self.ppt_title_box = None
        self.ppt_prs = None
        self._init_ppt_deck()
        
        self.tools_container = ctk.CTkFrame(self, fg_color="transparent")
        self.tools_container.grid(row=0, column=0, sticky="ew", padx=10, pady=5)
        
        # Row 1: Filters
        filter_frame = ctk.CTkFrame(self.tools_container, fg_color="transparent")
        filter_frame.pack(fill="x", pady=2)
        ctk.CTkLabel(filter_frame, text="Data Filter:", font=ctk.CTkFont(weight="bold")).pack(side="left", padx=5)
        
        filters = ["None", "Median Filter", "Moving Average Filter", "Kalman Filter"]
        if has_scipy:
            filters.extend([
                "Savitzky–Golay Filter", 
                "Butterworth Low-Pass Filter",
                "Elliptic Filter",
                "Hampel Filter",
                "Gaussian Filter",
                "Wiener Filter",
                "Notch Filter",
                "Wavelet Denoising Filter"
            ])
            
        self.combo_filter = ctk.CTkComboBox(filter_frame, values=filters, width=220)
        self.combo_filter.pack(side="left", padx=5)
        
        ctk.CTkLabel(filter_frame, text="Window/Cutoff/Freq/Sigma:").pack(side="left", padx=5)
        self.ent_filter_window = ctk.CTkEntry(filter_frame, width=60)
        self.ent_filter_window.insert(0, "10")
        self.ent_filter_window.pack(side="left", padx=5)
        
        btn_apply_filter = ctk.CTkButton(filter_frame, text="Apply Filter", width=100, command=self.draw_graph)
        btn_apply_filter.pack(side="left", padx=20)
        
        # Row 1.5: Sorting Row
        sort_frame = ctk.CTkFrame(self.tools_container, fg_color="transparent")
        sort_frame.pack(fill="x", pady=2)
        ctk.CTkLabel(sort_frame, text="Sort Data:", font=ctk.CTkFont(weight="bold")).pack(side="left", padx=5)
        
        self.combo_sort_col = ctk.CTkComboBox(sort_frame, values=["None"], width=150)
        self.combo_sort_col.pack(side="left", padx=5)
        
        self.combo_sort_order = ctk.CTkComboBox(sort_frame, values=["Ascending", "Descending"], width=120)
        self.combo_sort_order.pack(side="left", padx=5)
        
        btn_apply_sort = ctk.CTkButton(sort_frame, text="Apply Sort", width=100, command=self.draw_graph)
        btn_apply_sort.pack(side="left", padx=20)
        
        # Row 2: Texts Row
        text_frame = ctk.CTkFrame(self.tools_container, fg_color="transparent")
        text_frame.pack(fill="x", pady=2)
        ctk.CTkLabel(text_frame, text="Graph Texts:", font=ctk.CTkFont(weight="bold")).pack(side="left", padx=5)
        
        ctk.CTkLabel(text_frame, text="Title:").pack(side="left", padx=2)
        self.ent_title = ctk.CTkEntry(text_frame, width=120)
        self.ent_title.pack(side="left")
        
        ctk.CTkLabel(text_frame, text="X Label:").pack(side="left", padx=2)
        self.ent_xlabel = ctk.CTkEntry(text_frame, width=100)
        self.ent_xlabel.pack(side="left")
        
        ctk.CTkLabel(text_frame, text="Y1 Label:").pack(side="left", padx=2)
        self.ent_y1label = ctk.CTkEntry(text_frame, width=100)
        self.ent_y1label.pack(side="left")
        
        ctk.CTkLabel(text_frame, text="Y2 Label:").pack(side="left", padx=2)
        self.ent_y2label = ctk.CTkEntry(text_frame, width=80)
        self.ent_y2label.pack(side="left")
        
        ctk.CTkLabel(text_frame, text="Legend Pos:").pack(side="left", padx=2)
        self.combo_legend_loc = ctk.CTkComboBox(text_frame, values=["best", "upper right", "upper left", "lower left", "lower right", "right", "center"], width=100)
        self.combo_legend_loc.pack(side="left")
        
        # Row 3: Limits Row
        limit_frame = ctk.CTkFrame(self.tools_container, fg_color="transparent")
        limit_frame.pack(fill="x", pady=2)
        ctk.CTkLabel(limit_frame, text="Axes Limits:", font=ctk.CTkFont(weight="bold")).pack(side="left", padx=5)
        ctk.CTkLabel(limit_frame, text="X Min:").pack(side="left", padx=5)
        self.ent_xmin = ctk.CTkEntry(limit_frame, width=60)
        self.ent_xmin.pack(side="left")
        
        ctk.CTkLabel(limit_frame, text="X Max:").pack(side="left", padx=5)
        self.ent_xmax = ctk.CTkEntry(limit_frame, width=60)
        self.ent_xmax.pack(side="left")
        
        btn_y_limits = ctk.CTkButton(
            limit_frame,
            text="Y Limits",
            width=90,
            fg_color="#334155",
            hover_color="#1e293b",
            command=self.open_y_limits_dialog,
        )
        btn_y_limits.pack(side="left", padx=(18, 8))

        self.lbl_y_limits_status = ctk.CTkLabel(limit_frame, text="Auto Y", text_color="#64748b", width=62, anchor="w")
        self.lbl_y_limits_status.pack(side="left", padx=(0, 8))

        btn_apply = ctk.CTkButton(limit_frame, text="Apply", width=80, command=self.apply_limits)
        btn_apply.pack(side="left", padx=(0, 10))

        btn_clear = ctk.CTkButton(limit_frame, text="Clear Limits", width=100, fg_color="#F57C00", hover_color="#EF6C00", command=self.clear_limits)
        btn_clear.pack(side="left")
        
        self.var_show_grid = ctk.BooleanVar(value=True)
        chk_grid = ctk.CTkCheckBox(limit_frame, text="Show Grid", variable=self.var_show_grid, command=self.draw_graph)
        chk_grid.pack(side="left", padx=20)
        
        self.is_fullscreen = False
        btn_fs = ctk.CTkButton(limit_frame, text="⛶ Full Screen", width=120, fg_color="#2563EB", hover_color="#1E40AF", command=self.toggle_fullscreen)
        btn_fs.pack(side="left", padx=10)
        
        # Matplotlib Figure
        self.fig, self.ax = plt.subplots(figsize=(8, 5))
        self.fig.tight_layout(pad=3.0)
        self.ax2 = None
        self.y_axes = [self.ax]
        
        self.graph_content_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.graph_content_frame.grid(row=4, column=0, sticky="nsew", padx=10, pady=5)
        self.graph_content_frame.grid_rowconfigure(0, weight=1)
        self.graph_content_frame.grid_columnconfigure(0, weight=1)
        self.graph_content_frame.grid_columnconfigure(1, weight=0)

        self.canvas = FigureCanvasTkAgg(self.fig, master=self.graph_content_frame)
        self.canvas.get_tk_widget().grid(row=0, column=0, sticky="nsew", padx=(0, 10), pady=0)
        
        self.stats_box = ctk.CTkTextbox(self.graph_content_frame, height=140, font=("Consolas", 12), fg_color="transparent")
        self.stats_box.grid(row=1, column=0, sticky="ew", padx=(0, 10), pady=(5, 0))
        self.stats_box.insert("1.0", "Graph Statistics will appear here.")
        self.stats_box.configure(state="disabled")

        self._build_ppt_preview_panel()
        
        self._panning = False
        self._pan_start = None

        def _is_over_legend(event):
            for graph_ax in self._iter_y_axes():
                leg = graph_ax.get_legend()
                if leg and leg.contains(event)[0]:
                    return True
            return False

        def on_press(event):
            if event.button == 1 and event.inaxes and not _is_over_legend(event) and self.toolbar.mode == "":
                self._panning = True
                self._pan_start = (event.x, event.y)

        def on_drag(event):
            if not self._panning or not event.inaxes: return
            dx = event.x - self._pan_start[0]
            dy = event.y - self._pan_start[1]
            self._pan_start = (event.x, event.y)
            
            ax = event.inaxes
            x_min, x_max = self.ax.get_xlim()
            
            x_disp = (x_max - x_min) * (dx / ax.bbox.width)
            
            self.ax.set_xlim([x_min - x_disp, x_max - x_disp])
            for graph_ax in self._iter_y_axes():
                y_min, y_max = graph_ax.get_ylim()
                y_disp = (y_max - y_min) * (dy / graph_ax.bbox.height)
                graph_ax.set_ylim([y_min - y_disp, y_max - y_disp])
                
            self.canvas.draw_idle()

        def on_release(event):
            if event.button == 1:
                self._panning = False
                
        def zoom_fun(event):
            ax = event.inaxes
            if not ax: return
            base_scale = 1.2
            scale_factor = 1/base_scale if event.button == 'up' else base_scale
            xdata, ydata = event.xdata, event.ydata
            xlim = self.ax.get_xlim()
            self.ax.set_xlim([xdata - (xdata-xlim[0])*scale_factor, xdata + (xlim[1]-xdata)*scale_factor])
            for graph_ax in self._iter_y_axes():
                ylim = graph_ax.get_ylim()
                yrange = ylim[1] - ylim[0]
                ymid = ylim[0] + yrange / 2
                if graph_ax is ax and ydata is not None:
                    graph_ax.set_ylim([ydata - (ydata-ylim[0])*scale_factor, ydata + (ylim[1]-ydata)*scale_factor])
                else:
                    graph_ax.set_ylim([ymid - (yrange*scale_factor)/2, ymid + (yrange*scale_factor)/2])
            self.canvas.draw_idle()
            
        self.fig.canvas.mpl_connect('scroll_event', zoom_fun)
        self.fig.canvas.mpl_connect('button_press_event', on_press)
        self.fig.canvas.mpl_connect('motion_notify_event', on_drag)
        self.fig.canvas.mpl_connect('button_release_event', on_release)
        
        # Add Context Menu
        self.context_menu = tk.Menu(self, tearoff=0)
        self.context_menu.add_command(label="Copy Image to Clipboard", command=self.copy_image)
        self.context_menu.add_command(label="Export Editable Chart to Excel", command=self.export_editable)
        
        self.canvas.get_tk_widget().bind("<Button-3>", self.show_context_menu)
        self.canvas.get_tk_widget().bind("<Button-2>", self.show_context_menu)
        
        # Toolbar
        self.toolbar_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.toolbar_frame.grid(row=1, column=0, sticky="ew", padx=10)
        self.toolbar = NavigationToolbar2Tk(self.canvas, self.toolbar_frame)
        self.toolbar.update()
        
        self.btn_exit_fs = ctk.CTkButton(self.toolbar_frame, text="Exit Full Screen", width=120, fg_color="#F44336", hover_color="#D32F2F", command=self.toggle_fullscreen)
        import sys
        if sys.platform != 'darwin':
            self.controller.bind("<Escape>", lambda e: self.exit_fs_esc())

    def resource_path(self, filename):
        if getattr(sys, "frozen", False):
            base_dir = getattr(sys, "_MEIPASS", os.path.dirname(sys.executable))
        else:
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        return os.path.join(base_dir, filename)

    def _iter_y_axes(self):
        axes = [self.ax]
        axes.extend([axis for axis in getattr(self, "extra_axes", []) if axis is not None])
        return axes

    def _sync_y_axis_limit_values(self):
        for y_col, entries in self.y_axis_limit_entries.items():
            self.y_axis_limit_values[y_col] = {
                "min": entries["min"].get().strip(),
                "max": entries["max"].get().strip(),
            }

    def _refresh_y_limit_controls(self, y_cols):
        self._sync_y_axis_limit_values()
        self.y_axis_limit_entries = {}
        self._update_y_limits_status(y_cols)

    def _update_y_limits_status(self, y_cols=None):
        if y_cols is None:
            y_cols = getattr(self.controller, "graph_config", {}).get("y_cols", [])

        active_count = 0
        for y_col in y_cols:
            limits = self.y_axis_limit_values.get(y_col, {})
            if limits.get("min") or limits.get("max"):
                active_count += 1

        status = "Auto Y" if active_count == 0 else f"{active_count} set"
        if hasattr(self, "lbl_y_limits_status"):
            self.lbl_y_limits_status.configure(text=status)

    def open_y_limits_dialog(self):
        config = getattr(self.controller, "graph_config", {})
        y_cols = config.get("y_cols", [])
        if not y_cols:
            messagebox.showwarning("Y Limits", "Please configure graph Y-axis columns first.", parent=self)
            return

        self._sync_y_axis_limit_values()

        dialog = ctk.CTkToplevel(self)
        dialog.title("Y-Axis Limits")
        dialog.geometry("560x420")
        dialog.transient(self)
        dialog.grab_set()

        header = ctk.CTkFrame(dialog, fg_color="transparent")
        header.pack(fill="x", padx=20, pady=(18, 8))
        ctk.CTkLabel(header, text="Y-Axis Limits", font=ctk.CTkFont(size=18, weight="bold")).pack(anchor="w")
        ctk.CTkLabel(header, text="Leave min or max blank to keep that side on auto scale.", text_color="#64748b").pack(anchor="w", pady=(2, 0))

        rows = ctk.CTkScrollableFrame(dialog, fg_color="#f8fafc", corner_radius=8, height=250)
        rows.pack(fill="both", expand=True, padx=20, pady=(4, 12))

        dialog_entries = {}
        for idx, y_col in enumerate(y_cols, start=1):
            values = self.y_axis_limit_values.get(y_col, {"min": "", "max": ""})
            row = ctk.CTkFrame(rows, fg_color="#ffffff", corner_radius=6)
            row.pack(fill="x", padx=6, pady=5)

            label = ctk.CTkLabel(row, text=f"Y{idx}: {y_col}", width=210, anchor="w")
            label.pack(side="left", padx=(10, 8), pady=8)
            ctk.CTkLabel(row, text="Min").pack(side="left", padx=(0, 4))
            ent_min = ctk.CTkEntry(row, width=86)
            ent_min.insert(0, values.get("min", ""))
            ent_min.pack(side="left", padx=(0, 10))
            ctk.CTkLabel(row, text="Max").pack(side="left", padx=(0, 4))
            ent_max = ctk.CTkEntry(row, width=86)
            ent_max.insert(0, values.get("max", ""))
            ent_max.pack(side="left", padx=(0, 10))
            dialog_entries[y_col] = {"min": ent_min, "max": ent_max}

        footer = ctk.CTkFrame(dialog, fg_color="transparent")
        footer.pack(fill="x", padx=20, pady=(0, 18))

        def collect_dialog_values():
            for y_col, entries in dialog_entries.items():
                self.y_axis_limit_values[y_col] = {
                    "min": entries["min"].get().strip(),
                    "max": entries["max"].get().strip(),
                }

        def clear_y_values():
            for entries in dialog_entries.values():
                entries["min"].delete(0, "end")
                entries["max"].delete(0, "end")

        def apply_and_close():
            collect_dialog_values()
            self._update_y_limits_status(y_cols)
            self.apply_limits(redraw=True)
            dialog.destroy()

        ctk.CTkButton(footer, text="Clear Y Limits", width=120, fg_color="#64748b", hover_color="#475569", command=clear_y_values).pack(side="left")
        ctk.CTkButton(footer, text="Cancel", width=100, fg_color="#64748b", hover_color="#475569", command=dialog.destroy).pack(side="right")
        ctk.CTkButton(footer, text="Apply", width=110, command=apply_and_close).pack(side="right", padx=(0, 10))
        dialog.bind("<Return>", lambda _event: apply_and_close())

    def _build_y_axes(self, y_cols):
        self.ax2 = None
        self.extra_axes = []
        self.y_axes = [self.ax]

        if self._use_shared_db_axis(y_cols):
            self.y_axes = [self.ax] * max(1, len(y_cols))
            self.fig.subplots_adjust(right=0.84)
            return

        for idx in range(1, len(y_cols)):
            axis = self.ax.twinx()
            if idx > 1:
                axis.spines["right"].set_position(("axes", 1 + 0.08 * (idx - 1)))
                axis.set_frame_on(True)
                axis.patch.set_visible(False)
            self.extra_axes.append(axis)
            self.y_axes.append(axis)

        self.ax2 = self.y_axes[1] if len(self.y_axes) > 1 else None
        right_margin = 0.84 if len(y_cols) <= 2 else max(0.62, 0.84 - (len(y_cols) - 2) * 0.055)
        self.fig.subplots_adjust(right=right_margin)

    def _use_shared_db_axis(self, y_cols):
        return bool(y_cols) and all(str(col).strip().endswith("dB(A)") for col in y_cols)

    def _collect_legend_items(self):
        handles = []
        labels = []
        for graph_ax in self._iter_y_axes():
            axis_handles, axis_labels = graph_ax.get_legend_handles_labels()
            handles.extend(axis_handles)
            labels.extend(axis_labels)
        return handles, labels

    def _init_ppt_deck(self):
        if not has_pptx or not os.path.exists(self.ppt_template_path):
            return

        try:
            self.ppt_prs = Presentation(self.ppt_template_path)
            template_slide = self.ppt_prs.slides[0]
            for shape in template_slide.shapes:
                shape_name = str(shape.name).strip().lower()
                if shape_name == "graph":
                    self.ppt_graph_box = (shape.left, shape.top, shape.width, shape.height)
                if shape_name == "title" or (
                    shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
                ):
                    self.ppt_title_box = (shape.left, shape.top, shape.width, shape.height)

            if self.ppt_graph_box is None:
                graph_height = Cm(14.04)
                graph_width = Cm(25.4)
                self.ppt_graph_box = (Emu(0), self.ppt_prs.slide_height - graph_height, graph_width, graph_height)
        except Exception:
            self.ppt_prs = None

    def _build_ppt_preview_panel(self):
        self.ppt_panel = ctk.CTkFrame(self.graph_content_frame, width=285, corner_radius=8, fg_color="#ffffff")
        self.ppt_panel.grid(row=0, column=1, sticky="ns", padx=(0, 0), pady=0)
        self.ppt_panel.grid_propagate(False)

        header = ctk.CTkFrame(self.ppt_panel, fg_color="transparent")
        header.pack(fill="x", padx=12, pady=(12, 6))
        ctk.CTkLabel(header, text="PPT Preview", font=ctk.CTkFont(size=16, weight="bold")).pack(side="left")
        self.lbl_ppt_count = ctk.CTkLabel(header, text="0 slides", text_color="#64748b")
        self.lbl_ppt_count.pack(side="right")

        action_row = ctk.CTkFrame(self.ppt_panel, fg_color="transparent")
        action_row.pack(fill="x", padx=12, pady=(0, 8))
        self.btn_add_ppt = ctk.CTkButton(action_row, text="Add+", width=76, command=self.open_add_slide_dialog)
        self.btn_add_ppt.pack(side="left")
        self.btn_download_ppt = ctk.CTkButton(
            action_row,
            text="Download PPT",
            width=128,
            fg_color="#0f766e",
            hover_color="#115e59",
            command=self.download_ppt,
            state="disabled",
        )
        self.btn_download_ppt.pack(side="right")

        if not has_pptx or self.ppt_prs is None:
            self.btn_add_ppt.configure(state="disabled")
            self.btn_download_ppt.configure(state="disabled")
            message = "PPT export needs python-pptx and Temp.pptx."
        else:
            message = "Added slides will appear here."

        self.ppt_preview_scroll = ctk.CTkScrollableFrame(self.ppt_panel, fg_color="#f8fafc", corner_radius=6)
        self.ppt_preview_scroll.pack(fill="both", expand=True, padx=12, pady=(0, 12))
        self.lbl_ppt_empty = ctk.CTkLabel(self.ppt_preview_scroll, text=message, text_color="#64748b", wraplength=225)
        self.lbl_ppt_empty.pack(padx=12, pady=28)

    def open_add_slide_dialog(self):
        if self.ppt_prs is None:
            messagebox.showerror("PPT Export", "Temp.pptx could not be loaded.", parent=self)
            return
        if not self._has_plotted_graph():
            messagebox.showwarning("PPT Export", "Please generate a graph before adding a slide.", parent=self)
            return

        dialog = ctk.CTkToplevel(self)
        dialog.title("Add Graph to PPT")
        dialog.geometry("420x210")
        dialog.transient(self)
        dialog.grab_set()

        ctk.CTkLabel(dialog, text="Slide Title", font=ctk.CTkFont(size=16, weight="bold")).pack(anchor="w", padx=20, pady=(20, 6))
        title_entry = ctk.CTkEntry(dialog, placeholder_text="Enter title for this slide")
        title_entry.pack(fill="x", padx=20, pady=(0, 12))
        default_title = self.ent_title.get().strip()
        if default_title:
            title_entry.insert(0, default_title)
        title_entry.focus_set()

        footer = ctk.CTkFrame(dialog, fg_color="transparent")
        footer.pack(fill="x", padx=20, pady=(10, 20))

        def add_slide():
            title = title_entry.get().strip()
            if not title:
                messagebox.showwarning("Slide Title", "Please enter a title.", parent=dialog)
                return
            if self.add_current_graph_to_ppt(title):
                dialog.destroy()

        ctk.CTkButton(footer, text="Cancel", width=100, fg_color="#64748b", hover_color="#475569", command=dialog.destroy).pack(side="right")
        ctk.CTkButton(footer, text="Add Slide", width=120, command=add_slide).pack(side="right", padx=(0, 10))
        dialog.bind("<Return>", lambda _event: add_slide())

    def _has_plotted_graph(self):
        return any(len(axis.lines) > 0 for axis in self._iter_y_axes())

    def add_current_graph_to_ppt(self, title):
        try:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            image_path = os.path.join(self.ppt_output_dir, f"graph_{timestamp}.png")
            self.fig.savefig(image_path, format="png", dpi=300, bbox_inches="tight")

            slide = self._get_ppt_slide_for_add()
            self._fill_slide_title(slide, title)
            self._replace_slide_graph(slide, image_path)

            self.ppt_slide_entries.append({"title": title, "image_path": image_path})
            self._refresh_ppt_preview()
            messagebox.showinfo("PPT Export", "Graph added to PPT preview.", parent=self)
            return True
        except Exception as e:
            messagebox.showerror("PPT Export", f"Failed to add graph to PPT:\n{e}", parent=self)
            return False

    def _get_ppt_slide_for_add(self):
        if len(self.ppt_slide_entries) == 0:
            return self.ppt_prs.slides[0]
        return self.ppt_prs.slides.add_slide(self.ppt_prs.slides[0].slide_layout)

    def _fill_slide_title(self, slide, title):
        title_shape = None
        for shape in slide.shapes:
            if str(shape.name).strip().lower() == "title":
                title_shape = shape
                break
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                title_shape = shape
                break

        if title_shape is None:
            left, top, width, height = self.ppt_title_box or (Cm(1.5), Cm(0.3), Cm(22), Cm(2))
            title_shape = slide.shapes.add_textbox(left, top, width, height)
            title_shape.name = "Title"

        title_shape.name = "Title"
        title_shape.text = title

    def _replace_slide_graph(self, slide, image_path):
        for shape in list(slide.shapes):
            is_title = shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
            is_graph = str(shape.name).strip().lower() == "graph"
            if is_graph or (shape.is_placeholder and not is_title):
                element = shape._element
                element.getparent().remove(element)

        left, top, width, height = self.ppt_graph_box
        picture = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        picture.name = "graph"

    def _refresh_ppt_preview(self):
        for widget in self.ppt_preview_scroll.winfo_children():
            widget.destroy()

        self.ppt_preview_images = []
        slide_count = len(self.ppt_slide_entries)
        self.lbl_ppt_count.configure(text=f"{slide_count} slide{'s' if slide_count != 1 else ''}")
        self.btn_download_ppt.configure(state="normal" if slide_count else "disabled")

        if not slide_count:
            self.lbl_ppt_empty = ctk.CTkLabel(
                self.ppt_preview_scroll,
                text="Added slides will appear here.",
                text_color="#64748b",
                wraplength=250,
            )
            self.lbl_ppt_empty.pack(padx=12, pady=28)
            return

        for idx, entry in enumerate(self.ppt_slide_entries, start=1):
            card = ctk.CTkFrame(self.ppt_preview_scroll, fg_color="#ffffff", corner_radius=6)
            card.pack(fill="x", padx=4, pady=6)

            top_row = ctk.CTkFrame(card, fg_color="transparent")
            top_row.pack(fill="x", padx=10, pady=(8, 4))
            ctk.CTkLabel(top_row, text=f"Slide {idx}", font=ctk.CTkFont(weight="bold")).pack(side="left")
            ctk.CTkLabel(top_row, text=entry["title"], text_color="#334155", wraplength=165, justify="right").pack(side="right")

            try:
                img = Image.open(entry["image_path"]).convert("RGB")
                img.thumbnail((230, 126))
                ctk_img = ctk.CTkImage(light_image=img, dark_image=img, size=img.size)
                self.ppt_preview_images.append(ctk_img)
                ctk.CTkLabel(card, text="", image=ctk_img).pack(padx=10, pady=(0, 10))
            except Exception:
                ctk.CTkLabel(card, text="Preview unavailable", text_color="#64748b").pack(padx=10, pady=(0, 10))

    def download_ppt(self):
        if self.ppt_prs is None or not self.ppt_slide_entries:
            messagebox.showwarning("PPT Export", "No PPT slides are ready to download.", parent=self)
            return

        save_path = filedialog.asksaveasfilename(
            parent=self,
            title="Save PPT",
            defaultextension=".pptx",
            filetypes=[("PowerPoint Presentation", "*.pptx")],
            initialfile=f"SRDP_Graphs_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
        )
        if not save_path:
            return

        try:
            self.ppt_prs.save(save_path)
            messagebox.showinfo("PPT Export", f"PPT saved to:\n{save_path}", parent=self)
        except Exception as e:
            messagebox.showerror("PPT Export", f"Failed to save PPT:\n{e}", parent=self)

    def exit_fs_esc(self):
        if self.is_fullscreen: self.toggle_fullscreen()

    def toggle_fullscreen(self):
        if not self.is_fullscreen:
            # Enter Full Screen
            self.controller.top_bar.grid_forget()
            self.controller.sidebar_frame.grid_forget()
            self.controller.main_frame.grid(row=0, column=0, rowspan=2, columnspan=2, sticky="nsew")
            self.tools_container.grid_forget()
            self.ppt_panel.grid_forget()
            self.canvas.get_tk_widget().grid_configure(padx=0)
            self.btn_exit_fs.pack(side="right", padx=10)
            self.controller.state('zoomed')
            self.is_fullscreen = True
        else:
            # Exit Full Screen
            self.controller.top_bar.grid(row=0, column=0, columnspan=2, sticky="new")
            self.controller.sidebar_frame.grid(row=1, column=0, sticky="nsew")
            self.controller.main_frame.grid(row=1, column=1, rowspan=1, columnspan=1, sticky="nsew")
            self.tools_container.grid(row=0, column=0, sticky="ew", padx=10, pady=5)
            self.ppt_panel.grid(row=0, column=1, sticky="ns", padx=(0, 0), pady=0)
            self.canvas.get_tk_widget().grid_configure(padx=(0, 10))
            self.btn_exit_fs.pack_forget()
            self.is_fullscreen = False

    def show_context_menu(self, event):
        self.context_menu.tk_popup(event.x_root, event.y_root)

    def copy_image(self):
        try:
            import io
            import win32clipboard
            from PIL import Image
            buf = io.BytesIO()
            self.fig.savefig(buf, format='png', bbox_inches='tight', dpi=300)
            buf.seek(0)
            image = Image.open(buf)
            output = io.BytesIO()
            image.convert("RGB").save(output, "BMP")
            data = output.getvalue()[14:]
            
            win32clipboard.OpenClipboard()
            win32clipboard.EmptyClipboard()
            win32clipboard.SetClipboardData(win32clipboard.CF_DIB, data)
            win32clipboard.CloseClipboard()
            messagebox.showinfo("Success", "High-res graph image copied to clipboard!")
        except Exception as e:
            messagebox.showerror("Error", f"Clipboard error: {e}")

    def export_editable(self):
        from tkinter import filedialog
        save_path = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel Files", "*.xlsx")], initialfile="Exported_Chart.xlsx")
        if not save_path: return
        
        try:
            from openpyxl import Workbook
            from openpyxl.chart import ScatterChart, Reference, Series
            from openpyxl.chart.axis import ChartLines
            
            wb = Workbook()
            ws = wb.active
            ws.title = "Graph Data"
            
            chart = ScatterChart()
            chart.title = self.ent_title.get() or "Exported Graph"
            chart.x_axis.title = self.ent_xlabel.get() or "X Axis"
            chart.y_axis.title = self.ent_y1label.get() or "Y Axis"
            chart.style = 2
            chart.scatterStyle = "line"
            
            # Explicitly force numbers and labels to show
            chart.x_axis.delete = False
            chart.y_axis.delete = False
            chart.x_axis.tickLblPos = "low"
            chart.y_axis.tickLblPos = "low"
            
            # Use 'min' crossing to avoid hiding numbers inside data range
            chart.x_axis.crosses = "min"
            chart.y_axis.crosses = "min"
            
            # Style the axes
            chart.x_axis.majorTickMark = "out"
            chart.y_axis.majorTickMark = "out"
            
            try:
                xmin, xmax = self.ax.get_xlim()
            except:
                xmin, xmax = -np.inf, np.inf
            
            col_offset = 1
            for ax_idx, ax in enumerate(self.fig.axes):
                for line in ax.lines:
                    xdata = np.array(line.get_xdata())
                    ydata = np.array(line.get_ydata())
                    label = str(line.get_label())
                    
                    if len(xdata) == 0: continue
                    
                    mask = (xdata >= min(xmin, xmax)) & (xdata <= max(xmin, xmax))
                    x_clipped = xdata[mask]
                    y_clipped = ydata[mask]
                    
                    if len(x_clipped) == 0: continue
                    
                    ws.cell(row=1, column=col_offset, value=f"{label} X")
                    ws.cell(row=1, column=col_offset+1, value=f"{label} Y")
                    
                    for r, (xv, yv) in enumerate(zip(x_clipped, y_clipped), start=2):
                        ws.cell(row=r, column=col_offset, value=float(xv))
                        ws.cell(row=r, column=col_offset+1, value=float(yv))

                    x_ref = Reference(ws, min_col=col_offset, min_row=2, max_row=len(x_clipped)+1)
                    y_ref = Reference(ws, min_col=col_offset+1, min_row=2, max_row=len(y_clipped)+1)
                    
                    series_label = label if (label and label != '_nolegend_') else f"Line {col_offset}"
                    series = Series(y_ref, x_ref, title=series_label)
                    series.marker.symbol = "none"
                    
                    try:
                        line_w = float(self.controller.settings_manager.settings.get("line_width", 1.5))
                        series.graphicalProperties.line.width = int(line_w * 12700) # EMUs
                        
                        mpl_color = line.get_color()
                        if isinstance(mpl_color, str) and mpl_color.startswith('#'):
                            hex_raw = mpl_color.lstrip('#').upper()
                        else:
                            hex_raw = mcolors.to_hex(mpl_color).lstrip('#').upper()
                        series.graphicalProperties.line.solidFill = hex_raw
                    except:
                        pass
                        
                    chart.series.append(series)
                    col_offset += 2
            
            chart.legend.position = "r"
            chart.width = 30
            chart.height = 15
            ws.add_chart(chart, "D4")
            wb.save(save_path)
            messagebox.showinfo("Success", f"Editable MS Excel Chart saved to {save_path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to export: {e}")

    def on_show(self):
        if not hasattr(self.controller, 'graph_config') or not self.controller.graph_config.get("x_col"):
            self.ax.clear()
            for axis in getattr(self, "extra_axes", []):
                try:
                    axis.remove()
                except ValueError:
                    pass
            self.extra_axes = []
            self.ax2 = None
            self.y_axes = [self.ax]
            self._refresh_y_limit_controls([])
            self.ax.text(0.5, 0.5, "Configure Graph First", ha='center', va='center', fontsize=20)
            self.canvas.draw()
            return
            
        self.draw_graph()

    def _build_xtrp_noise_summary(self, y_cols):
        xtrp_files = [
            file_info for file_info in self.controller.data_manager.files_data
            if file_info.get("source_format") == ".xtrp"
        ]
        if len(xtrp_files) < 2:
            return []

        selected_mics = {
            str(col).replace(" dB(A)", "").strip()
            for col in y_cols
            if str(col).strip().endswith("dB(A)")
        }
        all_rows = []
        by_mic = {}
        for file_info in xtrp_files:
            details = file_info.get("import_details", {})
            metrics = details.get("xtrp_noise_metrics", {})
            if not metrics:
                continue

            tag = str(file_info.get("tag") or os.path.basename(file_info.get("original_filepath") or file_info.get("filepath", "XTRP")))
            for mic_name, metric in metrics.items():
                if selected_mics and mic_name not in selected_mics:
                    continue

                row = {
                    "mic": mic_name,
                    "tag": tag,
                    "overall": metric.get("overall_db_a", float("nan")),
                    "p95": metric.get("p95_db_a", float("nan")),
                    "peak": metric.get("peak_db_a", float("nan")),
                    "peak_frequency": metric.get("peak_frequency_hz", float("nan")),
                    "sharp_peaks": int(metric.get("sharp_peak_count", 0)),
                    "score": metric.get("score", float("inf")),
                }
                all_rows.append(row)
                by_mic.setdefault(mic_name, []).append(row)

        if len(all_rows) < 2:
            return []

        def fmt_db(value):
            return "n/a" if not np.isfinite(value) else f"{value:.2f}"

        def fmt_hz(value):
            return "n/a" if not np.isfinite(value) else f"{value:.1f}"

        def table(headers, rows):
            widths = [
                max(len(str(header)), *(len(str(row[idx])) for row in rows))
                for idx, header in enumerate(headers)
            ]
            result = [" | ".join(str(header).ljust(widths[idx]) for idx, header in enumerate(headers))]
            result.append("-+-".join("-" * width for width in widths))
            for row in rows:
                result.append(" | ".join(str(value).ljust(widths[idx]) for idx, value in enumerate(row)))
            return result

        detail_rows = []
        for row in sorted(all_rows, key=lambda item: (item["mic"], item["score"], item["tag"])):
            detail_rows.append((
                row["mic"],
                row["tag"],
                fmt_db(row["overall"]),
                fmt_db(row["p95"]),
                f"{fmt_db(row['peak'])} @ {fmt_hz(row['peak_frequency'])}",
                row["sharp_peaks"],
            ))

        verdict_rows = []
        for mic_name, rows in sorted(by_mic.items()):
            if len(rows) < 2:
                continue

            ranked = sorted(rows, key=lambda item: (item["score"], item["overall"], item["peak"], item["sharp_peaks"]))
            best = ranked[0]
            next_best = ranked[1]
            rms_advantage = next_best["overall"] - best["overall"]
            peak_advantage = next_best["peak"] - best["peak"]
            peak_note = "no sharp peaks" if best["sharp_peaks"] == 0 else f"{best['sharp_peaks']} sharp peaks"
            if best["sharp_peaks"] < next_best["sharp_peaks"]:
                peak_note += ", fewer than next"
            elif best["sharp_peaks"] == next_best["sharp_peaks"]:
                peak_note += ", same as next"
            else:
                peak_note += ", more than next"

            verdict_rows.append((
                mic_name,
                best["tag"],
                f"{fmt_db(rms_advantage)} dB lower",
                f"{fmt_db(peak_advantage)} dB peak gap",
                peak_note,
            ))

        lines = ["", "XTRP A-Weighted Frequency Summary (Frequency Hz vs Sound Level dB(A)):"]
        lines.extend(table(
            ["Mic Position", "File", "Overall A-RMS dB(A)", "95% dB(A)", "Peak dB(A) @ Hz", "Sharp Peaks"],
            detail_rows,
        ))

        if verdict_rows:
            lines.extend(["", "Least Noise / Smoothest Verdict by Mic Position:"])
            lines.extend(table(
                ["Mic Position", "Best File", "RMS Advantage", "Peak Advantage", "Peak Smoothness"],
                verdict_rows,
            ))

            overall_scores = {}
            for row in all_rows:
                overall_scores.setdefault(row["tag"], []).append(row["score"])
            overall_best = min(overall_scores.items(), key=lambda item: np.nanmean(item[1]))[0]
            lines.append(
                f"Overall Summary: {overall_best} is the lowest-noise/smoothest file across the selected mic positions "
                "by the combined A-weighted RMS, high-level content, peak level, and sharp-peak score."
            )
        return lines
        
    def draw_graph(self):
        self._sync_y_axis_limit_values()
        self.ax.clear()
        for axis in getattr(self, "extra_axes", []):
            try:
                axis.remove()
            except ValueError:
                pass
        self.extra_axes = []
        self.ax2 = None
        
        config = self.controller.graph_config
        x_col = config["x_col"]
        y_cols = config["y_cols"]
        self._refresh_y_limit_controls(y_cols)
        
        sort_values = ["None", x_col] + y_cols
        current_sort = self.combo_sort_col.get()
        self.combo_sort_col.configure(values=sort_values)
        if current_sort not in sort_values:
            self.combo_sort_col.set("None")
        
        self._build_y_axes(y_cols)
            
        filter_type = self.combo_filter.get()
        try:
            window_val = float(self.ent_filter_window.get())
        except ValueError:
            window_val = 10.0
            
        self._plot_counter = 0
        try:
            custom_colors = self.controller.settings_manager.settings.get("line_colors", [])
        except:
            custom_colors = []

        plotted_any = False
        categorical_ticks = None
        stats_lines = []
        
        sort_col = self.combo_sort_col.get()
        ascending = (self.combo_sort_order.get() == "Ascending")
        
        for file_info in self.controller.data_manager.files_data:
            df = file_info["df"].copy()
            tag = file_info["tag"]
            
            if sort_col != "None" and sort_col in df.columns:
                try:
                    df = df.sort_values(by=sort_col, ascending=ascending, key=lambda col: pd.to_numeric(col, errors='ignore'))
                except Exception:
                    df = df.sort_values(by=sort_col, ascending=ascending)
            
            if x_col in df.columns:
                x_raw = df[x_col]
                x_data = pd.to_numeric(x_raw, errors='coerce')
                x_non_empty = x_raw.notna() & (x_raw.astype(str).str.strip() != "")
                use_numeric_x = x_data.notna().sum() >= max(2, int(x_non_empty.sum() * 0.6))
                
                for idx, y_col in enumerate(y_cols):
                    if y_col in df.columns:
                        y_data = pd.to_numeric(df[y_col], errors='coerce')

                        if use_numeric_x:
                            valid_idx = ~(x_data.isna() | y_data.isna())
                            x_clean = x_data[valid_idx]
                            y_clean = y_data[valid_idx]
                        else:
                            labels = x_raw.astype(str).str.strip()
                            valid_idx = x_non_empty & ~y_data.isna()
                            labels = labels[valid_idx].reset_index(drop=True)
                            y_clean = y_data[valid_idx].reset_index(drop=True)
                            x_clean = np.arange(len(labels))
                            if categorical_ticks is None and len(labels) > 0:
                                categorical_ticks = list(labels)
                        
                        if len(x_clean) == 0: continue
                        
                        if filter_type == "Median Filter" and window_val >= 1:
                            y_clean = pd.Series(apply_median(y_clean.values, window_val))
                        elif filter_type == "Moving Average Filter" and window_val >= 1:
                            y_clean = y_clean.rolling(window=int(window_val), min_periods=1).mean()
                        elif filter_type == "Savitzky–Golay Filter" and has_scipy:
                            y_clean = pd.Series(apply_savgol(y_clean.values, window=window_val))
                        elif filter_type == "Butterworth Low-Pass Filter" and has_scipy:
                            y_clean = pd.Series(apply_lpf(y_clean.values, cutoff=window_val))
                        elif filter_type == "Elliptic Filter" and has_scipy:
                            # Clamp explicitly between 1 to 50
                            clamped_cutoff = max(1.0, min(50.0, float(window_val)))
                            y_clean = pd.Series(apply_elliptic(y_clean.values, cutoff=clamped_cutoff, order=5))
                        elif filter_type == "Hampel Filter":
                            y_clean = pd.Series(apply_hampel(y_clean.values, window=window_val))
                        elif filter_type == "Gaussian Filter":
                            y_clean = pd.Series(apply_gaussian(y_clean.values, sigma=window_val))
                        elif filter_type == "Wiener Filter":
                            y_clean = pd.Series(apply_wiener(y_clean.values, window=window_val))
                        elif filter_type == "Kalman Filter":
                            m_var = window_val / 100.0 if window_val > 0 else 0.1
                            y_clean = pd.Series(apply_kalman(y_clean.values, measurement_var=m_var))
                        elif filter_type == "Notch Filter" and has_scipy:
                            y_clean = pd.Series(apply_notch(y_clean.values, freq=window_val))
                        elif filter_type == "Wavelet Denoising Filter":
                            y_clean = pd.Series(apply_wavelet(y_clean.values))
                            
                        try:
                            line_w = float(self.controller.settings_manager.settings.get("line_width", 1.5))
                        except ValueError:
                            line_w = 1.5
                            
                        target_ax = self.y_axes[idx] if idx < len(self.y_axes) else self.ax
                        
                        target_color = custom_colors[self._plot_counter] if self._plot_counter < len(custom_colors) and custom_colors[self._plot_counter] else None
                        
                        if len(y_cols) > 1:
                            series_label = f"{tag} - {y_col}"
                        else:
                            series_label = str(tag)
                            
                        if target_color:
                            target_ax.plot(x_clean, y_clean, label=series_label, linewidth=line_w, color=target_color)
                        else:
                            target_ax.plot(x_clean, y_clean, label=series_label, linewidth=line_w)
                            
                        try:
                            y_vals = np.array(y_clean)
                            x_vals = np.array(x_clean) if use_numeric_x else np.array(labels)
                            max_idx = np.argmax(y_vals)
                            min_idx = np.argmin(y_vals)
                            y_max = y_vals[max_idx]
                            y_min = y_vals[min_idx]
                            x_at_max = x_vals[max_idx]
                            x_at_min = x_vals[min_idx]
                            stats_lines.append(f"[{series_label}] Min Y: {y_min:.4g} (at X: {x_at_min}) | Max Y: {y_max:.4g} (at X: {x_at_max})")
                        except Exception:
                            pass
                            
                        plotted_any = True
                        self._plot_counter += 1

        if categorical_ticks:
            tick_count = len(categorical_ticks)
            if tick_count <= 40:
                tick_positions = np.arange(tick_count)
                tick_labels = categorical_ticks
            else:
                step = max(1, int(np.ceil(tick_count / 40)))
                tick_positions = np.arange(0, tick_count, step)
                tick_labels = [categorical_ticks[i] for i in tick_positions]

            self.ax.set_xticks(tick_positions)
            self.ax.set_xticklabels(tick_labels, rotation=45, ha="right")

        self.apply_texts()

        if not plotted_any:
            self.ax.text(
                0.5,
                0.5,
                "No plottable data found for the selected columns.",
                ha="center",
                va="center",
                fontsize=14,
                transform=self.ax.transAxes,
            )
        
        # Apply themes to both axes
        axes = self._iter_y_axes()
        
        try:
            loc = self.combo_legend_loc.get()
            if not loc: loc = "best"
        except:
            loc = "best"
            
        show_grid = self.var_show_grid.get()
        
        if self.controller.settings_manager.settings.get("theme") == "Dark":
            self.fig.patch.set_facecolor('#2b2b2b')
            for a in axes:
                a.set_facecolor('#2b2b2b')
                a.tick_params(colors='white')
                a.xaxis.label.set_color('white')
                a.yaxis.label.set_color('white')
                for spine in a.spines.values(): spine.set_edgecolor('white')
                
            if show_grid:
                self.ax.grid(True, color='#555555')
            else:
                self.ax.grid(False)
            
            handles, labels = self._collect_legend_items()
            if handles:
                leg = self.ax.legend(handles, labels, loc=loc, facecolor='#2b2b2b', edgecolor='white')
                if leg: leg.set_draggable(True)
                for text in leg.get_texts(): text.set_color('white')
        else:
            self.fig.patch.set_facecolor('white')
            for a in axes:
                a.set_facecolor('white')
                a.tick_params(colors='black')
                a.xaxis.label.set_color('black')
                a.yaxis.label.set_color('black')
                for spine in a.spines.values(): spine.set_edgecolor('black')
                
            if show_grid:
                self.ax.grid(True, color='#dddddd')
            else:
                self.ax.grid(False)
            
            handles, labels = self._collect_legend_items()
            if handles:
                leg = self.ax.legend(handles, labels, loc=loc)
                if leg: leg.set_draggable(True)
            
        self.stats_box.configure(state="normal")
        self.stats_box.delete("1.0", "end")
        xtrp_summary = self._build_xtrp_noise_summary(y_cols)
        if xtrp_summary:
            stats_lines.extend(xtrp_summary)
        if stats_lines:
            self.stats_box.insert("1.0", "\n".join(stats_lines))
        else:
            self.stats_box.insert("1.0", "No data plotted.")
        self.stats_box.configure(state="disabled")

        self.apply_limits(redraw=False)
        self.canvas.draw()
        
    def apply_texts(self):
        user_title = self.ent_title.get()
        user_xlabel = self.ent_xlabel.get()
        user_y1label = self.ent_y1label.get()
        
        config = self.controller.graph_config
        y_cols = config.get("y_cols", ["Y Axis"])
        x_col = config.get("x_col", "X Axis")
        
        if user_title: self.ax.set_title(user_title)
        self.ax.set_xlabel(user_xlabel if user_xlabel else x_col)
        default_y1label = "Sound Level dB(A)" if self._use_shared_db_axis(y_cols) else y_cols[0]
        self.ax.set_ylabel(user_y1label if user_y1label else default_y1label)

        if self._use_shared_db_axis(y_cols):
            return
        
        for idx, graph_ax in enumerate(self.y_axes[1:], start=1):
            if idx >= len(y_cols):
                continue
            if idx == 1:
                user_y2label = self.ent_y2label.get()
                graph_ax.set_ylabel(user_y2label if user_y2label else y_cols[idx])
            else:
                graph_ax.set_ylabel(y_cols[idx])
            
    def apply_limits(self, redraw=True):
        self.apply_texts()
        try:
            if self.ent_xmin.get(): self.ax.set_xlim(left=float(self.ent_xmin.get()))
            if self.ent_xmax.get(): self.ax.set_xlim(right=float(self.ent_xmax.get()))
        except ValueError:
            pass

        self._sync_y_axis_limit_values()
        config = getattr(self.controller, "graph_config", {})
        y_cols = config.get("y_cols", [])
        for idx, y_col in enumerate(y_cols):
            if idx >= len(self.y_axes):
                continue
            limits = self.y_axis_limit_values.get(y_col, {})
            try:
                y_min = limits.get("min", "")
                y_max = limits.get("max", "")
                if y_min:
                    self.y_axes[idx].set_ylim(bottom=float(y_min))
                if y_max:
                    self.y_axes[idx].set_ylim(top=float(y_max))
            except ValueError:
                continue

        self._update_y_limits_status(y_cols)
        
        if redraw:
            self.canvas.draw()

    def clear_limits(self):
        self.ent_xmin.delete(0, 'end')
        self.ent_xmax.delete(0, 'end')
        self.y_axis_limit_values = {}
        for entries in self.y_axis_limit_entries.values():
            entries["min"].delete(0, 'end')
            entries["max"].delete(0, 'end')
        self._update_y_limits_status()
        self.draw_graph()
